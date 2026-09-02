"""
final_rag/ingestion/parser.py

PDF/DOCX/PPTX/MD document parser for the RAG pipeline. Uses Gemini for
OCR/transcription of PDFs, and lightweight extractors for office files.
"""

from __future__ import annotations

import hashlib
import logging
import re
import shutil
import tempfile
import time
import uuid
from concurrent.futures import ThreadPoolExecutor, TimeoutError as FuturesTimeoutError, as_completed
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Optional

import pymupdf as fitz
import warnings
from google import genai

try:
    from final_rag.config import (
        GEMINI_API_KEY,
        GEMINI_MODEL,
        SUPPORTED_EXTENSIONS,
        MAX_FILE_SIZE_MB,
        GEMINI_MAX_OUTPUT_TOKENS,
        GEMINI_TEMPERATURE,
        PAGES_PER_BATCH,
        MAX_ATTEMPTS_PER_BATCH,
        MAX_SPLIT_DEPTH,
        UPLOAD_POLL_INTERVAL_SEC,
        UPLOAD_POLL_MAX_ATTEMPTS,
        GEMINI_RETRY_SLEEP_BUSY_SEC,
        GEMINI_RETRY_SLEEP_QUOTA_SEC,
        PARSE_FOLDER_MAX_WORKERS,
        MD_OUTPUT_DIR,
    )
except ImportError:
    from ..config import (
        GEMINI_API_KEY,
        GEMINI_MODEL,
        SUPPORTED_EXTENSIONS,
        MAX_FILE_SIZE_MB,
        GEMINI_MAX_OUTPUT_TOKENS,
        GEMINI_TEMPERATURE,
        PAGES_PER_BATCH,
        MAX_ATTEMPTS_PER_BATCH,
        MAX_SPLIT_DEPTH,
        UPLOAD_POLL_INTERVAL_SEC,
        UPLOAD_POLL_MAX_ATTEMPTS,
        GEMINI_RETRY_SLEEP_BUSY_SEC,
        GEMINI_RETRY_SLEEP_QUOTA_SEC,
        PARSE_FOLDER_MAX_WORKERS,
        MD_OUTPUT_DIR,
    )

from .parser_metadata import extract_doc_metadata

warnings.filterwarnings("ignore")

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)-8s | %(name)s | %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
logger = logging.getLogger("ingestion.parser")


# ── Regex constants ───────────────────────────────────────────────────────────
_PAGE_MARKER_RE  = re.compile(r'^#{1,2}\s*Page\s+(\d+)\s*$', re.MULTILINE | re.IGNORECASE)
_HEADING_LINE_RE = re.compile(r'^###\s+(.*)$')
_TABLE_LINE_RE   = re.compile(r'^\s*\|.*\|\s*$')
_CODE_FENCE_RE   = re.compile(r'^```[a-zA-Z]*\s*$')


class ExtractionMethod(str, Enum):
    GEMINI  = "gemini_multimodal"
    OFFICE  = "office_extract"
    DIRECT  = "direct_markdown"
    SKIPPED = "skipped"
    ERROR   = "error"


@dataclass
class DocumentMeta:
    doc_id:                str
    file_name:             str
    file_path:             str
    file_type:             str
    file_size_kb:          float
    page_count:            int
    has_tables:            bool
    parse_success:         bool
    filename_tokens:       list[str] = field(default_factory=list)
    doc_year:              str       = ""
    doc_lang:              str       = "en"
    warnings:              list[str] = field(default_factory=list)


@dataclass
class PageRecord:
    page_no:    int
    page_label: str
    text:       str
    char_start: int = 0
    char_end:   int = 0


@dataclass
class BlockRecord:
    block_type: str
    content:    str
    page_no:    int
    page_label: str  = ""
    section:    str  = ""
    is_table:   bool = False


@dataclass
class ParseResult:
    file_name:           str
    file_type:           str
    method_used:         ExtractionMethod
    markdown:            str
    meta:                Optional[DocumentMeta] = None
    total_pages:         int                    = 0
    processing_time_sec: float                  = 0.0
    error:               Optional[str]          = None
    success:             bool                   = True
    warnings:            list[str]              = field(default_factory=list)
    pages:               list[PageRecord]       = field(default_factory=list)
    blocks:              list[BlockRecord]      = field(default_factory=list)
    doc_id:              str                    = ""

    @property
    def doc_year(self) -> str:
        return self.meta.doc_year if self.meta else ""


class _TruncatedResponseError(Exception):
    """Raised internally when Gemini's response hit the output-token ceiling."""
    def __init__(self, partial_text: str):
        super().__init__("Gemini response truncated (MAX_TOKENS)")
        self.partial_text = partial_text


_gemini_client = None


def _get_gemini_client():
    global _gemini_client
    if _gemini_client is None:
        if not GEMINI_API_KEY:
            raise RuntimeError(
                "GEMINI_API_KEY is not set. Add GEMINI_API_KEY=... to your .env file."
            )
        _gemini_client = genai.Client(api_key=GEMINI_API_KEY)
    return _gemini_client


def _upload_pdf(client, file_path: Path):
    logger.info("Uploading '%s' to Gemini...", file_path.name)

    unique = f"{file_path.name}-{uuid.uuid4().hex}"
    safe_name = hashlib.md5(unique.encode("utf-8")).hexdigest()[:16] + file_path.suffix
    tmp_dir   = Path(tempfile.gettempdir())
    safe_path = tmp_dir / safe_name
    shutil.copyfile(file_path, safe_path)

    try:
        gfile = client.files.upload(file=str(safe_path))
        for _ in range(UPLOAD_POLL_MAX_ATTEMPTS):
            info = client.files.get(name=gfile.name)
            if info.state.name == "ACTIVE":
                return gfile
            if info.state.name == "FAILED":
                raise RuntimeError(f"Gemini file processing failed for {file_path.name}")
            time.sleep(UPLOAD_POLL_INTERVAL_SEC)
        raise RuntimeError(f"Timed out waiting for Gemini to process {file_path.name}")
    finally:
        safe_path.unlink(missing_ok=True)


def _build_extraction_prompt(start: int, end: int) -> str:
    return f"""You are an expert Document Transcription and OCR specialist.

The attached PDF contains text in Hindi and/or English. Transcribe ONLY pages {start} through {end} (inclusive). Use this EXACT format for every page in that range, with no exceptions:

## Page <page_number>

<the full verbatim text content of that page, preserving paragraph breaks>

Any tables on the page must be written as GitHub-flavored Markdown tables, directly inline in the page's content, e.g.:

| Column A | Column B |
|---|---|
| value | value |

Any headings or titles on the page should be written as: ### <heading text>
This includes numbered section headings (e.g. "3.7 State Cooperative Development Committee (SCDC)"), not just standalone titles — if a line reads like a section heading (short, visually set apart from surrounding text, often numbered), mark it with ### even if it is followed immediately by body text or bullet points.

STRICT RULES:
1. Do NOT translate anything. Keep Hindi in Hindi (Devanagari script) and English in English, exactly as written.
2. Transcribe every word, sentence, and paragraph. Do not summarize or omit anything.
3. Preserve the original structure and reading order as closely as possible.
4. If a page is genuinely blank or has no visible content, output only:
   ## Page <page_number>

   *(blank page)*
5. NEVER write placeholder text such as "Screenshot for page N", "[image]", or "content not available". Always transcribe the actual visible content. If a word is unclear, give your best-effort reading rather than omitting it.
6. Output a "## Page <page_number>" marker for every single page in the requested range, in order.

Begin your transcription now, starting with page {start}.
"""


def _call_gemini(client, gfile, prompt: str) -> str:
    last_error = None

    for attempt in range(1, MAX_ATTEMPTS_PER_BATCH + 1):
        try:
            response = client.models.generate_content(
                model=GEMINI_MODEL,
                contents=[gfile, prompt],
                config={
                    "max_output_tokens": GEMINI_MAX_OUTPUT_TOKENS,
                    "temperature": GEMINI_TEMPERATURE,
                },
            )

            finish_reason = None
            try:
                finish_reason = response.candidates[0].finish_reason
            except Exception:
                pass

            text = response.text or ""

            if finish_reason is not None and "MAX_TOKENS" in str(finish_reason):
                raise _TruncatedResponseError(text)

            if not text.strip():
                raise RuntimeError("Gemini returned an empty response")

            return text

        except _TruncatedResponseError:
            raise

        except Exception as e:
            err_str = str(e)
            is_busy  = "503" in err_str or "UNAVAILABLE" in err_str
            is_quota = "429" in err_str or "RESOURCE_EXHAUSTED" in err_str

            if is_busy or is_quota:
                sleep_sec = GEMINI_RETRY_SLEEP_QUOTA_SEC if is_quota else GEMINI_RETRY_SLEEP_BUSY_SEC
                logger.warning(
                    "Gemini %s (attempt %d/%d). Sleeping %ds: %s",
                    "quota" if is_quota else "busy",
                    attempt, MAX_ATTEMPTS_PER_BATCH, sleep_sec, e,
                )
                time.sleep(sleep_sec)
                last_error = e
            else:
                logger.error("Gemini call failed: %s", e)
                raise

    raise RuntimeError(
        f"Gemini call failed after {MAX_ATTEMPTS_PER_BATCH} attempts. Last error: {last_error}"
    )


def _extract_page_range(client, gfile, start: int, end: int, depth: int = 0) -> str:
    prompt = _build_extraction_prompt(start, end)
    try:
        return _call_gemini(client, gfile, prompt)
    except _TruncatedResponseError as e:
        if depth >= MAX_SPLIT_DEPTH or start == end:
            logger.warning("Max split depth reached at page %d-%d. Keeping partial text.", start, end)
            return e.partial_text

        mid = (start + end) // 2
        logger.info("Batch %d-%d truncated. Halving -> [%d-%d] + [%d-%d]", start, end, start, mid, mid + 1, end)
        first_half  = _extract_page_range(client, gfile, start, mid, depth + 1)
        second_half = _extract_page_range(client, gfile, mid + 1, end, depth + 1)
        return first_half.rstrip() + "\n\n" + second_half.lstrip()


def _get_pdf_page_count(file_path: Path) -> int:
    doc = fitz.open(str(file_path))
    count = len(doc)
    doc.close()
    return count


def _extract_pdf_via_gemini(file_path: Path, total_pages: int, doc_warnings: list[str]) -> str:
    client = _get_gemini_client()
    gfile  = _upload_pdf(client, file_path)

    try:
        batches = [
            (b_start, min(b_start + PAGES_PER_BATCH - 1, total_pages))
            for b_start in range(1, total_pages + 1, PAGES_PER_BATCH)
        ]
        logger.info("Extracting %s | %d pages in %d batch(es)", file_path.name, total_pages, len(batches))

        chunks: list[str] = []
        for i, (b_start, b_end) in enumerate(batches, 1):
            logger.info("Batch %d/%d (pages %d-%d)...", i, len(batches), b_start, b_end)
            text = _extract_page_range(client, gfile, b_start, b_end)
            chunks.append(text)

        return "\n\n".join(chunks)

    finally:
        try:
            client.files.delete(name=gfile.name)
            logger.info("Deleted remote file '%s' from Gemini.", file_path.name)
        except Exception as e:
            logger.warning("Failed to delete remote file '%s': %s", file_path.name, e)


def _extract_office_document(file_path: Path, suffix: str, doc_warnings: list[str]) -> tuple[str, int]:
    if suffix == ".docx":
        import docx
        doc   = docx.Document(str(file_path))
        lines = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            lines.append("\n" + "\n".join(" | ".join(c.text.strip() for c in row.cells) for row in table.rows))
        return "## Page 1\n\n" + "\n\n".join(lines), 1

    if suffix == ".pptx":
        from pptx import Presentation
        prs   = Presentation(str(file_path))
        pages = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if shape.has_text_frame and shape.text.strip()
            ]
            pages.append(f"## Page {i}\n\n" + "\n\n".join(texts))
        return "\n\n".join(pages), len(prs.slides)

    raise ValueError(f"Unsupported office format: {suffix}")


def _split_into_pages(markdown: str) -> list[tuple[int, str]]:
    matches = list(_PAGE_MARKER_RE.finditer(markdown))
    if not matches:
        return [(1, markdown.strip())]

    pages: list[tuple[int, str]] = []
    for i, m in enumerate(matches):
        page_no     = int(m.group(1))
        content_start = m.end()
        content_end   = matches[i + 1].start() if i + 1 < len(matches) else len(markdown)
        page_text     = markdown[content_start:content_end].strip()
        pages.append((page_no, page_text))

    return pages


def _content_to_blocks(content: str, page_no: int, current_section: str) -> tuple[list[BlockRecord], str]:
    blocks: list[BlockRecord] = []
    lines  = content.splitlines()

    in_code      = False
    in_table     = False
    table_lines: list[str] = []
    text_lines:  list[str] = []

    def flush_text():
        if text_lines:
            txt = "\n".join(text_lines).strip()
            if txt:
                blocks.append(BlockRecord(
                    block_type="text",
                    content=txt,
                    page_no=page_no,
                    section=current_section,
                    is_table=False,
                ))
            text_lines.clear()

    for line in lines:
        if _CODE_FENCE_RE.match(line.strip()):
            in_code = not in_code
            text_lines.append(line)
            continue

        if not in_code:
            h_match = _HEADING_LINE_RE.match(line.strip())
            if h_match:
                flush_text()
                if in_table:
                    blocks.append(BlockRecord(
                        block_type="table",
                        content="\n".join(table_lines).strip(),
                        page_no=page_no,
                        section=current_section,
                        is_table=True,
                    ))
                    table_lines.clear()
                    in_table = False

                current_section = h_match.group(1).strip()
                blocks.append(BlockRecord(
                    block_type="heading",
                    content=current_section,
                    page_no=page_no,
                    section=current_section,
                    is_table=False,
                ))
                continue

            if _TABLE_LINE_RE.match(line):
                flush_text()
                in_table = True
                table_lines.append(line)
                continue
            elif in_table:
                blocks.append(BlockRecord(
                    block_type="table",
                    content="\n".join(table_lines).strip(),
                    page_no=page_no,
                    section=current_section,
                    is_table=True,
                ))
                table_lines.clear()
                in_table = False

        text_lines.append(line)

    flush_text()
    if in_table and table_lines:
        blocks.append(BlockRecord(
            block_type="table",
            content="\n".join(table_lines).strip(),
            page_no=page_no,
            section=current_section,
            is_table=True,
        ))

    return blocks, current_section


def _parse_markdown_document(markdown: str) -> tuple[list[PageRecord], list[BlockRecord]]:
    pages_raw = _split_into_pages(markdown)
    pages:  list[PageRecord]  = []
    blocks: list[BlockRecord] = []
    current_section = ""

    for page_no, content in pages_raw:
        pages.append(PageRecord(page_no=page_no, page_label=str(page_no), text=content))
        page_blocks, current_section = _content_to_blocks(content, page_no, current_section)
        blocks.extend(page_blocks)

    return pages, blocks


def _clean_noise(md: str) -> str:
    cleaned = []
    for line in md.splitlines():
        if re.search(r'^\s*page\s+\d+\s+of\s+\d+\s*$', line, re.IGNORECASE):
            continue
        cleaned.append(line)
    return "\n".join(cleaned)


def _postprocess_markdown(md: str) -> str:
    md = _clean_noise(md)
    md = re.sub(r'\n{3,}', '\n\n', md)
    return md.strip()


def _make_doc_id(file_name: str, file_size_kb: float) -> str:
    raw = f"{file_name}:{file_size_kb}"
    return hashlib.md5(raw.encode("utf-8")).hexdigest()[:12]


class DocumentParser:
    def __init__(self, output_dir: Optional[Path] = None):
        self.output_dir = Path(output_dir) if output_dir else MD_OUTPUT_DIR

    def parse_bytes(self, file_bytes: bytes, file_name: str, domain: str = "") -> ParseResult:
        suffix       = Path(file_name).suffix.lower()
        file_size_mb = len(file_bytes) / (1024 * 1024)

        if suffix not in SUPPORTED_EXTENSIONS:
            return self._skipped(file_name, suffix, f"Unsupported: {suffix}")
        if file_size_mb > MAX_FILE_SIZE_MB:
            return self._skipped(file_name, suffix, f"Too large: {file_size_mb:.1f}MB")

        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
            tmp.write(file_bytes)
            tmp_path = Path(tmp.name)
        try:
            return self.parse_file(tmp_path, original_file_name=file_name, domain=domain)
        finally:
            tmp_path.unlink(missing_ok=True)

    def parse_file(
        self,
        file_path:          str | Path,
        original_file_name: Optional[str] = None,
        domain:             str           = "",
    ) -> ParseResult:
        file_path = Path(file_path)
        start     = time.perf_counter()
        suffix    = file_path.suffix.lower()

        if not file_path.exists() or suffix not in SUPPORTED_EXTENSIONS:
            return self._skipped(file_path.name, suffix, f"Missing or unsupported: {file_path}")

        file_size_mb = file_path.stat().st_size / (1024 * 1024)
        file_size_kb = file_path.stat().st_size / 1024

        if file_size_mb > MAX_FILE_SIZE_MB:
            return self._skipped(file_path.name, suffix, f"Too large: {file_size_mb:.1f}MB")

        name_for_meta = original_file_name or file_path.name
        doc_id        = _make_doc_id(name_for_meta, round(file_size_kb, 2))

        logger.info("Parsing | %s | %.1f MB", name_for_meta, file_size_mb)
        doc_warnings: list[str] = []

        try:
            if suffix == ".pdf":
                total_pages  = _get_pdf_page_count(file_path)
                raw_markdown = _extract_pdf_via_gemini(file_path, total_pages, doc_warnings)
                method       = ExtractionMethod.GEMINI

            elif suffix == ".md":
                raw_markdown = file_path.read_text(encoding="utf-8", errors="replace")
                total_pages  = 1
                method       = ExtractionMethod.DIRECT

            elif suffix in (".docx", ".pptx"):
                raw_markdown, total_pages = _extract_office_document(file_path, suffix, doc_warnings)
                method = ExtractionMethod.OFFICE

            else:
                return self._skipped(name_for_meta, suffix, f"Unsupported: {suffix}")

        except Exception as e:
            logger.error("Parse failed | %s | %s", name_for_meta, e)
            return ParseResult(
                file_name           = name_for_meta,
                file_type           = suffix,
                method_used         = ExtractionMethod.ERROR,
                markdown            = "",
                success             = False,
                error               = str(e),
                warnings            = doc_warnings,
                processing_time_sec = round(time.perf_counter() - start, 3),
            )

        raw_markdown  = _postprocess_markdown(raw_markdown)
        pages, blocks = _parse_markdown_document(raw_markdown)
        self._save(file_path, raw_markdown, name_for_meta, doc_id, domain=domain)

        elapsed    = round(time.perf_counter() - start, 3)
        has_tables = any(b.is_table for b in blocks)

        doc_meta = extract_doc_metadata(name_for_meta, raw_markdown)

        meta = DocumentMeta(
            doc_id=doc_id,
            file_name=name_for_meta,
            file_path=str(file_path),
            file_type=suffix,
            file_size_kb=round(file_size_kb, 2),
            page_count=len(pages),
            has_tables=has_tables,
            parse_success=True,
            filename_tokens=doc_meta.filename_tokens,
            doc_year=doc_meta.doc_year,
            doc_lang=doc_meta.doc_lang,
            warnings=doc_warnings,
        )

        logger.info(
            "Done | %s | pages=%d | blocks=%d | tables=%d | lang=%s | year=%s | time=%.3fs",
            name_for_meta, len(pages), len(blocks),
            sum(1 for b in blocks if b.is_table),
            doc_meta.doc_lang, doc_meta.doc_year or "unknown", elapsed,
        )

        return ParseResult(
            file_name=name_for_meta,
            file_type=suffix,
            method_used=method,
            markdown=raw_markdown,
            meta=meta,
            total_pages=total_pages,
            success=True,
            warnings=doc_warnings,
            pages=pages,
            blocks=blocks,
            processing_time_sec=elapsed,
            doc_id=doc_id,
        )

    def parse_folder(
        self,
        folder_path: str | Path,
        recursive:   bool = False,
        max_workers: int  = PARSE_FOLDER_MAX_WORKERS,
    ) -> list[ParseResult]:
        folder_path = Path(folder_path)
        if not folder_path.exists() or not folder_path.is_dir():
            raise ValueError(f"Folder not found: {folder_path}")

        pattern = "**/*" if recursive else "*"
        files   = [
            f for f in folder_path.glob(pattern)
            if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS
        ]

        logger.info("Parsing %d file(s) | workers=%d", len(files), max_workers)

        if max_workers == 1:
            results = [self.parse_file(f) for f in files]
        else:
            results = []
            with ThreadPoolExecutor(max_workers=max_workers) as executor:
                futures = {executor.submit(self._parse_file_wrapper, f): f for f in files}
                for future in as_completed(futures):
                    fp = futures[future]
                    try:
                        results.append(future.result())
                    except FuturesTimeoutError:
                        logger.error("Timed out | %s", fp)
                        results.append(self._skipped(fp.name, fp.suffix, "Timed out"))
                    except Exception as e:
                        logger.error("Failed | %s | %s", fp, e)
                        results.append(self._skipped(fp.name, fp.suffix, str(e)))

        success = sum(1 for r in results if r.success)
        logger.info(
            "Batch done | total=%d | ok=%d | failed=%d",
            len(results), success, len(results) - success,
        )
        return results

    def _parse_file_wrapper(self, file_path: Path) -> ParseResult:
        try:
            return self.parse_file(file_path)
        except Exception as e:
            return self._skipped(file_path.name, file_path.suffix, str(e))

    def _save(
        self,
        file_path: Path,
        markdown: str,
        original_name: str = None,
        doc_id: str = "",
        domain: str = "",
    ) -> None:
        if self.output_dir and markdown.strip():
            stem = Path(original_name).stem if original_name else file_path.stem
            suffix = f"__{doc_id}" if doc_id else ""
            target_dir = (self.output_dir / domain) if domain else self.output_dir
            target_dir.mkdir(parents=True, exist_ok=True)
            (target_dir / f"{stem}{suffix}.md").write_text(markdown, encoding="utf-8")

    def _skipped(self, file_name: str, file_type: str, reason: str) -> ParseResult:
        logger.warning("Skipped | %s | %s", file_name, reason)
        return ParseResult(
            file_name   = file_name,
            file_type   = file_type,
            method_used = ExtractionMethod.SKIPPED,
            markdown    = "",
            success     = False,
            error       = reason,
        )